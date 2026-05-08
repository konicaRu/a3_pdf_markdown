from __future__ import annotations

import io
import re
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable

from PIL import Image, UnidentifiedImageError

from a3_pdf_markdown.app.converters.vision_client import VisionClient, build_vision_client
from a3_pdf_markdown.app.core.models import AppConfig, LogLevel


PPTX_IMAGE_PROMPT = (
    "Опиши изображение из презентации по-русски. Если это график или диаграмма, "
    "укажи тип, подписи, видимые значения и ключевой вывод. Если это схема, "
    "перечисли блоки и связи. Если это обычная иллюстрация, кратко опиши смысл. "
    "Не пересказывай весь слайд и не добавляй вводных фраз."
)


@dataclass(slots=True)
class PptxImageDescription:
    slide_number: int
    image_number: int | None
    description: str
    is_slide_level: bool = False


SLIDE_MARKER_RE = re.compile(r"(<!--\s*Slide number:\s*(\d+)\s*-->)")
IMAGE_MARKDOWN_RE = re.compile(r"(!\[[^\]]*]\([^)]+\))")


class PptxEnricher:
    def __init__(self, config: AppConfig, log: Callable[[LogLevel, str], None]) -> None:
        self.config = config
        self.log = log
        self._vision: VisionClient | None = None

    def describe_images(self, source_path: Path) -> list[PptxImageDescription]:
        if not self.config.vision_enabled:
            return []

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            self.log(LogLevel.WARNING, "Для описания картинок PPTX нужен пакет python-pptx")
            return []

        presentation = Presentation(str(source_path))
        descriptions: list[PptxImageDescription] = []
        slide_fallbacks: set[int] = set()

        for slide_index, slide in enumerate(presentation.slides, start=1):
            image_index = 0
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue

                image_index += 1
                try:
                    image_blob = shape.image.blob
                    image = Image.open(io.BytesIO(image_blob))
                    self.log(
                        LogLevel.INFO,
                        f"PPTX слайд {slide_index}: описание изображения {image_index}",
                    )
                    description = self._vision_client().describe_image(image, PPTX_IMAGE_PROMPT)
                except ValueError as exc:
                    if "no embedded image" in str(exc):
                        self.log(
                            LogLevel.WARNING,
                            f"PPTX слайд {slide_index}, изображение {image_index}: "
                            "нет embedded image, попробую описать слайд целиком",
                        )
                        slide_fallbacks.add(slide_index)
                        continue
                    raise
                except (UnidentifiedImageError, OSError) as exc:
                    self.log(
                        LogLevel.WARNING,
                        f"PPTX слайд {slide_index}, изображение {image_index}: "
                        f"не удалось прочитать изображение: {exc}",
                    )
                    continue
                except Exception as exc:
                    self.log(
                        LogLevel.WARNING,
                        f"PPTX слайд {slide_index}, изображение {image_index}: "
                        f"vision не смог описать изображение: {exc}",
                    )
                    continue

                if description.strip():
                    descriptions.append(
                        PptxImageDescription(
                            slide_number=slide_index,
                            image_number=image_index,
                            description=description.strip(),
                        )
                    )

        for slide_index in sorted(slide_fallbacks):
            description = self._describe_rendered_slide(source_path, slide_index)
            if description:
                descriptions.append(
                    PptxImageDescription(
                        slide_number=slide_index,
                        image_number=None,
                        description=description,
                        is_slide_level=True,
                    )
                )

        return descriptions

    def _describe_rendered_slide(self, source_path: Path, slide_number: int) -> str:
        try:
            slide_image = self._render_slide_with_powerpoint(source_path, slide_number)
        except Exception as exc:
            self.log(
                LogLevel.WARNING,
                f"PPTX слайд {slide_number}: не удалось отрендерить слайд через PowerPoint: {exc}",
            )
            return ""

        try:
            self.log(LogLevel.INFO, f"PPTX слайд {slide_number}: описание слайда целиком")
            return self._vision_client().describe_image(slide_image, PPTX_IMAGE_PROMPT).strip()
        except Exception as exc:
            self.log(
                LogLevel.WARNING,
                f"PPTX слайд {slide_number}: vision не смог описать слайд целиком: {exc}",
            )
            return ""

    def _render_slide_with_powerpoint(self, source_path: Path | str, slide_number: int) -> Image.Image:
        try:
            import pythoncom
            import win32com.client
        except ImportError as exc:
            raise RuntimeError("для рендера PPTX нужен pywin32") from exc

        source = str(Path(source_path).resolve())
        with tempfile.TemporaryDirectory(prefix="a3_pptx_slide_") as temp_dir:
            output_path = Path(temp_dir) / f"slide_{slide_number}.png"
            pythoncom.CoInitialize()
            app = None
            presentation = None
            try:
                app = win32com.client.DispatchEx("PowerPoint.Application")
                presentation = app.Presentations.Open(source, WithWindow=False)
                presentation.Slides(slide_number).Export(str(output_path), "PNG")
                with Image.open(output_path) as image:
                    return image.convert("RGB").copy()
            finally:
                if presentation is not None:
                    presentation.Close()
                if app is not None:
                    app.Quit()
                pythoncom.CoUninitialize()

    def append_descriptions(self, markdown: str, source_path: Path) -> tuple[str, bool]:
        descriptions = self.describe_images(source_path)
        if not descriptions:
            return markdown, False

        by_slide: dict[int, list[PptxImageDescription]] = {}
        for item in descriptions:
            by_slide.setdefault(item.slide_number, []).append(item)

        split = SLIDE_MARKER_RE.split(markdown)
        if len(split) == 1:
            return self._append_fallback(markdown, descriptions), True

        result = [split[0]]
        inserted_any = False
        for index in range(1, len(split), 3):
            marker = split[index]
            slide_number = int(split[index + 1])
            slide_body = split[index + 2]
            enriched_body, inserted = self._insert_into_slide(
                slide_body,
                by_slide.get(slide_number, []),
            )
            inserted_any = inserted_any or inserted
            result.extend([marker, enriched_body])

        if not inserted_any:
            return self._append_fallback(markdown, descriptions), True
        return "".join(result).strip() + "\n", True

    def _insert_into_slide(
        self,
        slide_body: str,
        descriptions: list[PptxImageDescription],
    ) -> tuple[str, bool]:
        if not descriptions:
            return slide_body, False

        by_image = {item.image_number: item for item in descriptions}
        image_counter = 0
        inserted = False

        def replace(match: re.Match[str]) -> str:
            nonlocal image_counter, inserted
            image_counter += 1
            item = by_image.get(image_counter)
            if item is None:
                return match.group(0)
            inserted = True
            return (
                f"{match.group(0)}\n\n"
                f"> **Описание изображения:** {item.description}\n"
            )

        enriched = IMAGE_MARKDOWN_RE.sub(replace, slide_body)
        missing = [
            item
            for item in descriptions
            if item.image_number is not None and item.image_number > image_counter
        ]
        slide_level = [item for item in descriptions if item.is_slide_level]
        if missing:
            inserted = True
            extra = ["\n\n### Описания изображений слайда\n"]
            for item in missing:
                extra.append(f"\n- Изображение {item.image_number}: {item.description}")
            enriched = enriched.rstrip() + "".join(extra) + "\n"
        if slide_level:
            inserted = True
            extra = ["\n\n### Описание визуальных элементов слайда\n"]
            for item in slide_level:
                extra.append(f"\n{item.description}\n")
            enriched = enriched.rstrip() + "".join(extra)
        return enriched, inserted

    def _append_fallback(
        self,
        markdown: str,
        descriptions: list[PptxImageDescription],
    ) -> str:
        parts = [markdown.rstrip(), "\n\n## Описания изображений\n"]
        for item in descriptions:
            label = "слайд целиком" if item.is_slide_level else f"изображение {item.image_number}"
            parts.append(
                f"\n### Слайд {item.slide_number}, {label}\n\n"
                f"{item.description}\n"
            )
        return "".join(parts).strip() + "\n"

    def _vision_client(self) -> VisionClient:
        if self._vision is None:
            self._vision = build_vision_client(self.config)
        return self._vision
